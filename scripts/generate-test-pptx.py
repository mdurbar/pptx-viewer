#!/usr/bin/env python3
"""
Generate test PPTX files for integration testing.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import nsmap
import os

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', 'tests', 'fixtures')

def ensure_output_dir():
    """Ensure the output directory exists."""
    os.makedirs(OUTPUT_DIR, exist_ok=True)

def create_basic_shapes_pptx():
    """Create a PPTX with various basic shapes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x41, 0x69, 0xE1)  # Royal Blue
    shape.text = "Rectangle"

    # Rounded Rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(0.5), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x32, 0xCD, 0x32)  # Lime Green
    shape.text = "Rounded Rect"

    # Ellipse
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(0.5), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0x69, 0xB4)  # Hot Pink
    shape.text = "Ellipse"

    # Triangle
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.5), Inches(2), Inches(2), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xA5, 0x00)  # Orange

    # Diamond
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(3), Inches(2), Inches(2), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x94, 0x00, 0xD3)  # Purple

    # Star
    shape = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, Inches(5.5), Inches(2), Inches(2), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xD7, 0x00)  # Gold

    # Arrow
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.5), Inches(4), Inches(2.5), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80)  # Teal

    # Heart
    shape = slide.shapes.add_shape(MSO_SHAPE.HEART, Inches(3.5), Inches(4), Inches(1.5), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)  # Red

    # Cloud
    shape = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(5.5), Inches(4), Inches(2.5), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x87, 0xCE, 0xEB)  # Sky Blue

    prs.save(os.path.join(OUTPUT_DIR, 'basic-shapes.pptx'))
    print("Created basic-shapes.pptx")

def create_text_formatting_pptx():
    """Create a PPTX with various text formatting options."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Text Formatting Test"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Bold, Italic, Underline
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]

    run = p.add_run()
    run.text = "Bold "
    run.font.bold = True
    run.font.size = Pt(18)

    run = p.add_run()
    run.text = "Italic "
    run.font.italic = True
    run.font.size = Pt(18)

    run = p.add_run()
    run.text = "Underline "
    run.font.underline = True
    run.font.size = Pt(18)

    run = p.add_run()
    run.text = "Strikethrough"
    run.font.size = Pt(18)
    # Note: strikethrough not directly supported in python-pptx

    # Different colors
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(0.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]

    colors = [
        ("Red ", RGBColor(0xFF, 0x00, 0x00)),
        ("Green ", RGBColor(0x00, 0x80, 0x00)),
        ("Blue ", RGBColor(0x00, 0x00, 0xFF)),
        ("Orange ", RGBColor(0xFF, 0xA5, 0x00)),
        ("Purple", RGBColor(0x80, 0x00, 0x80)),
    ]

    for text, color in colors:
        run = p.add_run()
        run.text = text
        run.font.color.rgb = color
        run.font.size = Pt(18)

    # Different sizes
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(8), Inches(1))
    tf = shape.text_frame
    p = tf.paragraphs[0]

    sizes = [12, 18, 24, 36, 48]
    for size in sizes:
        run = p.add_run()
        run.text = f"{size}pt "
        run.font.size = Pt(size)

    # Bullet list
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(4), Inches(2))
    tf = shape.text_frame
    tf.word_wrap = True

    items = ["First bullet item", "Second bullet item", "Third bullet item"]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)

    # Numbered list (simulated)
    shape = slide.shapes.add_textbox(Inches(5), Inches(3.8), Inches(4), Inches(2))
    tf = shape.text_frame
    tf.word_wrap = True

    items = ["1. First numbered item", "2. Second numbered item", "3. Third numbered item"]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)

    prs.save(os.path.join(OUTPUT_DIR, 'text-formatting.pptx'))
    print("Created text-formatting.pptx")

def create_gradients_pptx():
    """Create a PPTX with gradient fills."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Linear gradient (horizontal)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(3), Inches(2))
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 0
    fill.gradient_stops[0].color.rgb = RGBColor(0xFF, 0x00, 0x00)
    fill.gradient_stops[1].color.rgb = RGBColor(0x00, 0x00, 0xFF)
    shape.text = "Linear H"

    # Linear gradient (vertical)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(0.5), Inches(3), Inches(2))
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = RGBColor(0x00, 0xFF, 0x00)
    fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0xFF, 0x00)
    shape.text = "Linear V"

    # Diagonal gradient
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(3), Inches(2))
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = RGBColor(0x80, 0x00, 0x80)
    fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0xA5, 0x00)
    shape.text = "Diagonal"

    # Gradient on rounded rect
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(3), Inches(3), Inches(2))
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = RGBColor(0x00, 0x80, 0x80)
    fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0x69, 0xB4)
    shape.text = "Rounded Gradient"

    prs.save(os.path.join(OUTPUT_DIR, 'gradients.pptx'))
    print("Created gradients.pptx")

def create_images_pptx():
    """Create a PPTX with images (placeholder for now)."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Add a text note about images
    shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Image test - add images manually or use a sample image"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    prs.save(os.path.join(OUTPUT_DIR, 'images.pptx'))
    print("Created images.pptx (placeholder)")

def create_tables_pptx():
    """Create a PPTX with tables."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Table Test"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Simple table
    rows, cols = 4, 3
    left = Inches(1)
    top = Inches(1.2)
    width = Inches(8)
    height = Inches(2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for col in range(cols):
        table.columns[col].width = Inches(2.5)

    # Header row
    headers = ["Name", "Department", "Status"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x41, 0x69, 0xE1)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Data rows
    data = [
        ["Alice Smith", "Engineering", "Active"],
        ["Bob Johnson", "Marketing", "Active"],
        ["Carol Williams", "Sales", "On Leave"],
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value

    prs.save(os.path.join(OUTPUT_DIR, 'tables.pptx'))
    print("Created tables.pptx")

def create_multi_slide_pptx():
    """Create a PPTX with multiple slides."""
    prs = Presentation()

    # Slide 1 - Title
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Slide Presentation"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Testing slide navigation"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2 - Content
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Slide 2: Content Slide"
    p.font.size = Pt(32)
    p.font.bold = True

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x41, 0x69, 0xE1)
    shape.text = "Shape 1"

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(2), Inches(3), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x32, 0xCD, 0x32)
    shape.text = "Shape 2"

    # Slide 3 - More content
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Slide 3: Final Slide"
    p.font.size = Pt(32)
    p.font.bold = True

    shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "This is the final slide of the presentation. It contains text content to verify multi-slide navigation works correctly."
    p.font.size = Pt(18)

    prs.save(os.path.join(OUTPUT_DIR, 'multi-slide.pptx'))
    print("Created multi-slide.pptx")

def create_comprehensive_pptx():
    """Create a comprehensive test PPTX with many features."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Comprehensive Feature Test"
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

    # Row 1: Basic shapes with different fills
    shapes_data = [
        (MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1), Inches(1.8), Inches(1.2), RGBColor(0x41, 0x69, 0xE1), "Rect"),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(1), Inches(1.8), Inches(1.2), RGBColor(0x32, 0xCD, 0x32), "Round"),
        (MSO_SHAPE.OVAL, Inches(4.3), Inches(1), Inches(1.8), Inches(1.2), RGBColor(0xFF, 0x69, 0xB4), "Oval"),
        (MSO_SHAPE.DIAMOND, Inches(6.3), Inches(1), Inches(1.8), Inches(1.2), RGBColor(0xFF, 0xA5, 0x00), "Diamond"),
        (MSO_SHAPE.STAR_5_POINT, Inches(8.1), Inches(1), Inches(1.6), Inches(1.2), RGBColor(0xFF, 0xD7, 0x00), "Star"),
    ]

    for shape_type, left, top, width, height, color, text in shapes_data:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.text = text
        for p in shape.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    # Row 2: Arrows and more shapes
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.3), Inches(2.5), Inches(2), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x80)

    shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(2.5), Inches(2.5), Inches(2), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x80, 0x00, 0x80)

    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.7), Inches(2.5), Inches(2), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xDC, 0x14, 0x3C)

    shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(6.9), Inches(2.4), Inches(1.4), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x4B, 0x00, 0x82)

    shape = slide.shapes.add_shape(MSO_SHAPE.HEART, Inches(8.5), Inches(2.4), Inches(1.2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # Row 3: Text formatting samples
    shape = slide.shapes.add_textbox(Inches(0.3), Inches(3.6), Inches(4.5), Inches(1.2))
    tf = shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Bold "
    run.font.bold = True
    run.font.size = Pt(14)

    run = p.add_run()
    run.text = "Italic "
    run.font.italic = True
    run.font.size = Pt(14)

    run = p.add_run()
    run.text = "Underline "
    run.font.underline = True
    run.font.size = Pt(14)

    run = p.add_run()
    run.text = "Color"
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    run.font.size = Pt(14)

    # Gradient shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(3.6), Inches(4.5), Inches(1.2))
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = RGBColor(0x66, 0x00, 0xFF)
    fill.gradient_stops[1].color.rgb = RGBColor(0x00, 0xFF, 0xFF)
    shape.text = "Gradient Fill"
    for p in shape.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    # Row 4: Small table
    table = slide.shapes.add_table(2, 3, Inches(0.3), Inches(5.2), Inches(5), Inches(1)).table

    headers = ["Column A", "Column B", "Column C"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    data = ["Value 1", "Value 2", "Value 3"]
    for col, value in enumerate(data):
        cell = table.cell(1, col)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)

    # Text box with bullet points
    shape = slide.shapes.add_textbox(Inches(5.5), Inches(5.2), Inches(4), Inches(1.5))
    tf = shape.text_frame
    tf.word_wrap = True

    bullets = ["First item", "Second item", "Third item"]
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(12)

    prs.save(os.path.join(OUTPUT_DIR, 'comprehensive.pptx'))
    print("Created comprehensive.pptx")

def main():
    """Generate all test PPTX files."""
    ensure_output_dir()

    print(f"Generating test PPTX files in {OUTPUT_DIR}...")
    print("-" * 50)

    create_basic_shapes_pptx()
    create_text_formatting_pptx()
    create_gradients_pptx()
    create_tables_pptx()
    create_multi_slide_pptx()
    create_comprehensive_pptx()
    create_images_pptx()

    print("-" * 50)
    print("Done! Test files created in tests/fixtures/")

if __name__ == "__main__":
    main()
